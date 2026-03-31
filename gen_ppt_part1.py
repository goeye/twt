#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_slide(title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)
        if line.startswith('  '):
            p.level = 1
            p.font.size = Pt(16)

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
tf = title_box.text_frame
tf.text = "TWT Chat\n智能客服系统使用培训"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
tf = subtitle_box.text_frame
tf.text = "让客户服务更高效、更智能"
tf.paragraphs[0].font.size = Pt(24)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 2
add_slide("什么是 TWT Chat", [
    "企业级海外客服 SaaS 平台",
    "",
    "核心特点：",
    "• 多渠道统一管理",
    "  Email、Webhook 等多种渠道消息集中处理",
    "",
    "• AI 驱动的智能客服",
    "  基于知识库的自动回复，减少重复工作",
    "",
    "• 灵活的权限体系",
    "  支持自定义角色和精细化权限配置"
])

# Slide 3
add_slide("核心价值", [
    "为什么选择 TWT Chat？",
    "",
    "提升效率",
    "  AI 自动回复，减少 60% 重复工作",
    "",
    "多渠道整合",
    "  Email、Webhook 等统一管理，无需切换平台",
    "",
    "精准权限",
    "  灵活的角色与权限配置，保障数据安全",
    "",
    "数据驱动",
    "  完整的报表分析体系，持续优化服务"
])

# Slide 4
add_slide("功能架构", [
    "TWT Chat 五大核心模块",
    "",
    "• 客服工作台 - 会话管理、档案查询、客户管理",
    "",
    "• 渠道接入层 - Email、Webhook 多渠道集成",
    "",
    "• AI 智能引擎 - 自动回复、知识库驱动",
    "",
    "• 数据分析中心 - 报表统计、绩效分析",
    "",
    "• 权限管理系统 - 角色配置、坐席管理"
])

# Slide 5
add_slide("客服工作台 - 会话管理", [
    "实时会话处理中心",
    "",
    "实时会话处理",
    "  在线会话列表、排队提醒、状态更新",
    "",
    "会话操作",
    "  快速领取、转移会话、添加协作、标记/关闭",
    "",
    "辅助工具",
    "  快捷回复、标签分类、访客信息、历史记录",
    "",
    "工作流程：访客发起 → AI 回复 → 客服领取 → 处理 → 关闭"
])

prs.save('TWT-Chat-培训演示.pptx')
print("✅ PPT 已生成：TWT-Chat-培训演示.pptx")
