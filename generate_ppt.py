#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(content_lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)
        if line.startswith('•') or line.startswith('-') or line.startswith('✅') or line.startswith('📊'):
            p.level = 0
        elif line.startswith('  '):
            p.level = 1
            p.font.size = Pt(16)

    return slide

# Slide 1: 封面
add_title_slide("TWT Chat\n智能客服系统使用培训", "让客户服务更高效、更智能")

# Slide 2: 什么是 TWT Chat
add_content_slide("什么是 TWT Chat", [
    "TWT Chat - 企业级海外客服 SaaS 平台",
    "",
    "核心特点：",
    "• 🌐 多渠道统一管理",
    "  Email、Webhook 等多种渠道消息集中处理",
    "",
    "• 🤖 AI 驱动的智能客服",
    "  基于知识库的自动回复，减少重复工作",
    "",
    "• 🎯 灵活的权限体系",
    "  支持自定义角色和精细化权限配置",
    "",
    "适用场景：跨境电商、SaaS 企业、海外服务团队"
])

# Slide 3: 核心价值
add_content_slide("核心价值", [
    "为什么选择 TWT Chat？",
    "",
    "🚀 提升效率",
    "  AI 自动回复，减少 60% 重复工作",
    "",
    "🌐 多渠道整合",
    "  Email、Webhook 等统一管理，无需切换平台",
    "",
    "🎯 精准权限",
    "  灵活的角色与权限配置，保障数据安全",
    "",
    "📊 数据驱动",
    "  完整的报表分析体系，持续优化服务"
])

# Slide 4: 功能架构
add_content_slide("功能架构", [
    "TWT Chat 五大核心模块",
    "",
    "• 客服工作台 - 会话管理、档案查询、客户管理",
    "",
    "• 渠道接入层 - Email、Webhook 多渠道集成",
    "",
    "• AI 智能引擎 - 自动回复、知识库驱动",
    "",
    "• 数据分析中心 - 报表统计、绩效分析",
    "",
    "• 权限管理系统 - 角色配置、坐席管理"
])

# Slide 5: 客服工作台 - 会话管理
add_content_slide("客服工作台 - 会话管理", [
    "实时会话处理中心",
    "",
    "✅ 实时会话处理",
    "  在线会话列表、排队提醒、状态更新",
    "",
    "✅ 会话操作",
    "  快速领取、转移会话、添加协作、标记/关闭",
    "",
    "✅ 辅助工具",
    "  快捷回复、标签分类、访客信息、历史记录",
    "",
    "工作流程：访客发起 → AI 回复 → 客服领取 → 处理 → 关闭"
])

# Slide 6: 档案与客户
add_content_slide("客服工作台 - 档案与客户", [
    "完整的客户关系管理",
    "",
    "📁 档案模块",
    "  会话历史归档、多维度搜索、标签分类、质量分析",
    "",
    "👤 客户模块",
    "  客户资料管理、会话历史、标签体系、价值分析",
    "",
    "👁️ 访客模块",
    "  在线监控、主动营销、来源追踪、时长统计"
])

# Slide 7: AI Agent
add_content_slide("AI Agent - 智能客服", [
    "知识库驱动的自动化服务",
    "",
    "阶段一：访客首条消息",
    "  受众范围、回复时机、标签显示",
    "",
    "阶段二：AI 回复策略",
    "  严格/发散模式、知识库关联",
    "",
    "阶段三：无法解答处理",
    "  兜底回复、转接人工、提示语配置",
    "",
    "阶段四：访客跟进与自动关闭",
    "  不活跃跟进、自动关闭会话"
])

# Slide 8: 多渠道接入
add_content_slide("多渠道接入", [
    "统一管理，无缝集成",
    "",
    "✉️ Email 邮件会话系统",
    "  • 邮件线程智能识别（基于 Message-ID）",
    "  • 支持附件和富文本内容",
    "  • 换邮箱回复保持连续性",
    "",
    "🔗 Webhook 平台集成",
    "  • 支持：Slack、飞书、钉钉、企业微信、自定义",
    "  • HMAC-SHA256 签名验证",
    "  • 统一消息格式转换"
])

# Slide 9: 权限与角色管理
add_content_slide("权限与角色管理", [
    "灵活的权限体系",
    "",
    "系统预设角色",
    "  👑 管理员 - 全部权限，不可删除",
    "  👤 客服 - 基础权限，可编辑",
    "",
    "自定义角色：最多 99 个",
    "",
    "三级权限树",
    "  1. 菜单权限（12 个一级菜单）",
    "  2. 会话操作权限（查看、转移、关闭等）",
    "  3. 功能权限（快捷回复、标签管理等）"
])

# Slide 10: 营销与报表
add_content_slide("营销与报表", [
    "数据驱动的服务优化",
    "",
    "🎯 营销模块",
    "  主动营销活动、在线访客触达、产品推荐、活动推送",
    "",
    "📊 报表模块",
    "  • 客服绩效统计",
    "  • 响应时间分析",
    "  • 会话解决率",
    "  • 客户满意度",
    "  • 趋势对比分析",
    "",
    "价值：通过数据洞察，持续优化服务流程"
])

# Slide 11: 演示环节说明
add_content_slide("演示环节说明", [
    "实操演示准备",
    "",
    "🌐 演示环境",
    "  URL: https://pm.pro.jishu666.com/web-agent/",
    "  账号：管理员 / 客服",
    "",
    "演示流程",
    "  访客进入 → AI 回复 → 客服领取",
    "  → 快捷回复 → 添加标签 → 关闭会话",
    "",
    "演示内容",
    "  1. 客服接待完整流程",
    "  2. AI Agent 配置方法",
    "  3. 权限管理操作",
    "  4. 渠道配置步骤"
])

# Slide 12: 演示 1
add_content_slide("演示 1 - 客服接待流程", [
    "从访客进入到会话关闭",
    "",
    "1️⃣ 访客发起会话",
    "  访客在网站发送消息，进入排队队列",
    "",
    "2️⃣ AI 自动回复",
    "  AI 识别问题类型，基于知识库给出答案",
    "",
    "3️⃣ 客服领取会话",
    "  查看排队列表，点击"领取"按钮",
    "",
    "4️⃣ 使用快捷回复 → 5️⃣ 添加标签 → 6️⃣ 关闭会话"
])

# Slide 13: 演示 2
add_content_slide("演示 2 - AI Agent 配置", [
    "智能客服配置步骤",
    "",
    "步骤 1：进入 AI Agent 设置",
    "  导航：设置 → AI Agent",
    "",
    "步骤 2：配置受众范围",
    "  全部访客 / 仅访客 / 仅客户",
    "",
    "步骤 3：设置回复时机",
    "  始终回复 / 仅客服不在线时",
    "",
    "步骤 4-6：关联知识库、配置兜底回复、测试效果"
])

# Slide 14: 演示 3
add_content_slide("演示 3 - 权限管理", [
    "创建角色并分配权限",
    "",
    "1️⃣ 创建自定义角色",
    "  设置 → 客服管理 → 角色管理 → 新建角色",
    "",
    "2️⃣ 配置菜单权限",
    "  勾选可访问的功能模块",
    "",
    "3️⃣ 设置会话操作权限",
    "  勾选允许的操作（查看、转移、关闭等）",
    "",
    "4️⃣ 邀请客服成员",
    "  输入邮箱、选择角色、生成邀请链接"
])

# Slide 15: 演示 4
add_content_slide("演示 4 - 渠道配置", [
    "Email 和 Webhook 集成",
    "",
    "Email 渠道配置",
    "  1. 设置 → 邮箱渠道",
    "  2. 填写 SMTP/IMAP 服务器信息",
    "  3. 测试连接并保存",
    "",
    "Webhook 平台集成",
    "  1. 设置 → Webhook",
    "  2. 选择平台类型（Slack/飞书/钉钉等）",
    "  3. 填写 Webhook URL",
    "  4. 配置回调地址并测试"
])

# Slide 16: 使用建议
add_content_slide("使用建议", [
    "最佳实践分享",
    "",
    "1. 合理配置 AI Agent 受众范围",
    "  成熟客户群体 → AI 服务所有人",
    "",
    "2. 定期更新知识库内容",
    "  每周回顾聊天记录，补充高频问题",
    "",
    "3. 使用快捷回复提升效率",
    "  设置标准答案，一键发送",
    "",
    "4. 善用标签分类管理",
    "  方便统计分析和质量抽查",
    "",
    "5. 定期查看报表优化服务"
])

# Slide 17: 常见问题
add_content_slide("常见问题 Q&A", [
    "Q1: AI 无法回答时怎么办？",
    "💡 配置兜底回复 + 自动转接人工",
    "",
    "Q2: 如何管理多个客服团队？",
    "💡 使用自定义角色 + 精细化权限",
    "",
    "Q3: 邮件会话如何保持线程？",
    "💡 系统自动识别邮件头 Message-ID 字段",
    "",
    "Q4: 坐席名额不够怎么办？",
    "💡 联系管理员升级服务版本"
])

# Slide 18: 总结
add_content_slide("总结与答疑", [
    "TWT Chat 核心优势回顾",
    "",
    "🎯 四大核心价值",
    "  1. 多渠道统一管理 - 提升协作效率",
    "  2. AI 智能客服 - 减少重复工作",
    "  3. 灵活权限体系 - 适配不同团队",
    "  4. 完整数据分析 - 持续优化服务",
    "",
    "📞 联系方式",
    "  演示环境：https://pm.pro.jishu666.com/web-agent/",
    "",
    "🙋 开放问答环节 - 欢迎大家提问！"
])

print("PPT 生成中...")
prs.save('TWT-Chat-培训演示.pptx')
print("✅ PPT 已生成：TWT-Chat-培训演示.pptx")
